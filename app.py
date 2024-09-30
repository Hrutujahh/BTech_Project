from flask import Flask, render_template, request, send_file
import numpy as np
import cv2
import joblib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

app = Flask(__name__)
model = joblib.load('best_random_forest_model.pkl')

# Function to preprocess the image
def preprocess_image(img_path):
    img = cv2.imread(img_path)
    if img is not None:
        img = cv2.resize(img, (128, 128))  # Resize to 128x128
        img_flattened = img.flatten()
        return img_flattened
    return np.zeros((128 * 128 * 3))  # Fallback if image not found

# Function to predict features
def predict_features(thickness, diameter, force, material, img_path):
    input_data = np.array([[7850, 200000000000, 0.3, thickness, diameter, force, 0]])  # Max def placeholder
    predicted_features = model.predict(input_data)
    img = preprocess_image(img_path)
    return predicted_features, img

# Function to create a PowerPoint report
def create_report(image_path, predicted_features, thickness, diameter, force, material):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Stress and Deformation Report"
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(4.5), height=Inches(3))
    table = slide.shapes.add_table(8, 2, Inches(5), Inches(1.5), Inches(4), Inches(3)).table
    table.cell(0, 0).text = "Input Features"
    table.cell(0, 1).text = "Values"
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    input_data = [("Thickness (mm)", thickness), ("Diameter (mm)", diameter), ("Force (N)", force), ("Material", material), ("", ""), ("Predicted Features", "")]
    for i, (feature, value) in enumerate(input_data):
        table.cell(i + 1, 0).text = feature
        table.cell(i + 1, 1).text = str(value)
    max_stress, min_stress, max_deformation = predicted_features[0]
    predicted_data = [("Max Stress (MPa)", max_stress), ("Min Stress (MPa)", min_stress), ("Max Deformation (mm)", max_deformation)]
    for i, (feature, value) in enumerate(predicted_data):
        table.cell(i + 5, 0).text = feature
        table.cell(i + 5, 1).text = f"{value:.2f}"  # Format to 2 decimal places
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(12)  # Set font size for all cells
    report_path = 'static/stress_deformation_report.pptx'
    prs.save(report_path)
    return report_path

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'image' not in request.files:
        return 'No image uploaded', 400
    image_file = request.files['image']
    thickness = float(request.form['thickness'])
    diameter = float(request.form['diameter'])
    force = float(request.form['force'])
    material = request.form['material']

    # Save the uploaded image
    image_path = os.path.join('static', image_file.filename)
    image_file.save(image_path)

    # Make predictions
    predicted_features, _ = predict_features(thickness, diameter, force, material, image_path)

    # Create report
    report_path = create_report(image_path, predicted_features, thickness, diameter, force, material)

    return send_file(report_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
